"""
Shared helpers for reconstructing Figma slides as an editable PowerPoint deck.

Text layers become real, editable PowerPoint text boxes (font/size/color/position
matched to the Figma design, read via get_design_context). Everything else (photos,
icons, decorative shapes, rounded cards) is rendered as a flattened image exported
directly from Figma via get_screenshot on the specific non-text node -- this avoids
reverse-engineering Figma's CSS transforms (object-cover crops, mirrors, rotations)
and guarantees pixel-perfect backgrounds.

Two non-obvious python-pptx/OOXML fixes are baked into these helpers so every slide
built with them inherits both automatically (see memory: feedback_pptx_export_gotchas):
  - add_rounded_rect() disables the default theme shadow python-pptx puts on autoshapes.
  - add_text() converts line-height to exact Pt() spacing instead of a bare float, since
    a float puts PowerPoint into percentage mode (relative to its own font-metrics-based
    "single line spacing", taller than 1.0x font size) rather than matching Figma's
    line-height (an exact multiple of font size).
"""

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Figma canvas is 1920x1080px, treated as 144 "dpi" so 1920px -> 13.333in (16:9 widescreen).
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080

FONT_REGULAR = "Suisse Int'l"
FONT_MEDIUM = "Suisse Int'l Medium"
FONT_SEMIBOLD = "Suisse Int'l Semi Bold"


def px(v):
    return Emu(round(v * 6350))


def pt(figma_px):
    # 144 "dpi" canvas -> 2 Figma px per point
    return Pt(figma_px / 2)


def rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_corner_radius(shape, radius_px, box_w_px, box_h_px):
    """Apply a rounded-rectangle clip to a picture or autoshape via its adj value."""
    adj = min(0.5, radius_px / min(box_w_px, box_h_px))
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    geom = spPr.find(qn("a:prstGeom"))
    if geom is None:
        for tag in ("a:custGeom", "a:prstGeom"):
            existing = spPr.find(qn(tag))
            if existing is not None:
                spPr.remove(existing)
        geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
        spPr.append(geom)
    else:
        geom.set("prst", "roundRect")
        for child in list(geom):
            geom.remove(child)
    avLst = geom.makeelement(qn("a:avLst"), {})
    gd = avLst.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {int(adj * 100000)}"})
    avLst.append(gd)
    geom.append(avLst)


def add_picture(slide, path, x_px, y_px, w_px, h_px, radius_px=0):
    pic = slide.shapes.add_picture(path, px(x_px), px(y_px), px(w_px), px(h_px))
    if radius_px:
        set_corner_radius(pic, radius_px, w_px, h_px)
    return pic


def add_rounded_rect(slide, x_px, y_px, w_px, h_px, fill_hex, radius_px):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x_px), px(y_px), px(w_px), px(h_px))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.fill.background()
    shape.shadow.inherit = False
    set_corner_radius(shape, radius_px, w_px, h_px)
    return shape


def add_rect(slide, x_px, y_px, w_px, h_px, fill_hex):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x_px), px(y_px), px(w_px), px(h_px))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x_px, y_px, w_px, h_px, paragraphs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of paragraph specs, each a list of run tuples
    (text, font_name, size_px, color_hex, line_spacing) or a dict with 'runs' + 'space_after_px'."""
    box = slide.shapes.add_textbox(px(x_px), px(y_px), px(w_px), px(h_px))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, para_spec in enumerate(paragraphs):
        if isinstance(para_spec, dict):
            runs = para_spec["runs"]
            space_after = para_spec.get("space_after_px", 0)
        else:
            runs = para_spec
            space_after = 0

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        line_spacing = runs[0][4] if runs else None
        size_px = runs[0][2] if runs else None
        if line_spacing and size_px:
            p.line_spacing = Pt((size_px * line_spacing) / 2)
        if space_after:
            p.space_after = pt(space_after)
        else:
            p.space_after = Pt(0)
        p.space_before = Pt(0)

        for text, font_name, size_px, color_hex, _ls in runs:
            r = p.add_run()
            r.text = text
            r.font.name = font_name
            r.font.size = pt(size_px)
            r.font.color.rgb = rgb(color_hex)
        _prevent_orphan_word(p)
    return box


def _prevent_orphan_word(paragraph):
    """Join the last two words of a paragraph with a non-breaking space so word-wrap
    never leaves a single short word alone on the last line. A no-op if the text
    fits on one line anyway (NBSP renders identically to a regular space)."""
    for r in reversed(paragraph.runs):
        idx = r.text.rfind(" ")
        if idx != -1:
            r.text = r.text[:idx] + " " + r.text[idx + 1:]
            return


def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _approx_text_width_px(text, size_px):
    """Rough average-glyph-width estimate for Suisse Int'l, good enough for placing
    a trailing mark right after a short heading/label (not for precise line-wrapping)."""
    return len(text) * size_px * 0.56


def add_label_with_mark(slide, x, y, base_text, base_font, base_size_px, color_hex,
                         mark, mark_font, mark_size_px, w=900, h=70, line_spacing=1.1):
    """A heading/label followed by a trademark/registered mark, in its own separate
    text box top-aligned to the label -- rather than an inline run at a smaller size,
    which renders vertically centered mid-word instead of raised like a superscript."""
    add_text(slide, x, y, w, h, [[(base_text, base_font, base_size_px, color_hex, line_spacing)]])
    mark_x = x + _approx_text_width_px(base_text, base_size_px)
    add_text(slide, mark_x, y, 60, h, [[(mark, mark_font, mark_size_px, color_hex, line_spacing)]])


def add_meta(slide, left_text, right_text):
    add_text(slide, 48, 48, 300, 19, [[(left_text, FONT_REGULAR, 14, "#71717a", 1.0)]])
    add_text(slide, 1600, 48, 272, 19, [[(right_text, FONT_REGULAR, 14, "#71717a", 1.0)]], align=PP_ALIGN.RIGHT)
