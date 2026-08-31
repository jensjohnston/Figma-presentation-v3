"""Slide 01 -- Purifier Cover. Figma node 61219:17965 (canvas "All purifiers").

No standard meta-left/meta-right pair here -- only a single top-right line
("Bluewater Purifiers · 2026", darker #27272a than the usual gray500 meta
color), so add_meta() is skipped in favor of a direct add_text() call.

The three background product close-ups (Cleone, Spirit x2) are figma nodes
whose bounding boxes extend well past the 1920x1080 frame; get_screenshot
already returns them pre-clipped to the frame's visible region (Figma
respects the frame's own content-clip when rendering a descendant node), so
each screenshot's returned width/height is used as-is and positioned at
max(box_x, 0), max(box_y, 0).

Figma's own z-order draws these three photos *above* the title/meta/logo
layer, but only because their blank headspace is genuinely transparent
there -- get_screenshot flattens that onto an opaque matte instead
(feedback_pptx_export_gotchas #5), so copying the raw z-order here would
bury the title under an opaque near-white rectangle. Pictures are drawn
*first* instead, text last, so the title/meta/logo always render on top
regardless of the photo mattes.
"""

from pptx.enum.text import PP_ALIGN

from tools.pptx_core import new_slide, add_text, add_picture, add_rect, FONT_SEMIBOLD, FONT_REGULAR


def build(prs, assets_dir):
    slide = new_slide(prs)

    a = f"{assets_dir}/slide_01"

    add_rect(slide, 0, 0, 1920, 1080, "#f5f5f5")

    add_picture(slide, f"{a}/cleone_bg.png", 0, 0, 1561, 1080)
    add_picture(slide, f"{a}/spirit_bg_1.png", 782, 115, 1138, 965)
    add_picture(slide, f"{a}/spirit_bg_2.png", 0, 145, 1920, 935)

    add_text(
        slide, 1600, 48, 272, 19,
        [[("Bluewater Purifiers · 2026", FONT_REGULAR, 14, "#27272a", 1.0)]],
        align=PP_ALIGN.RIGHT,
    )
    add_picture(slide, f"{a}/logo_icon.png", 44, 48.5, 48, 46.737)
    add_text(
        slide, 48, 115, 1824, 150,
        [[("Bluewater Purifiers", FONT_SEMIBOLD, 130, "#18181b", 1.0)]],
        align=PP_ALIGN.CENTER,
    )

    return slide
