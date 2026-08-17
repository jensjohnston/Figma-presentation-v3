"""Slide 06 -- Spirit profile. Figma node 61219:18042.

Same card geometry as the main deck's slide_18_spirit.py (verified via
get_metadata: every card position matches exactly, including the clipped
510x204 fiber-thumb crop reused directly from that slide's asset pass) --
only the product photo, NSF badge and SuperiorOsmosis dome are freshly
screenshotted since those specific nodes/crops differ from the old slide.
A background card (595.92x920) now sits behind the product photo, matching
Cleone/Pro's composition instead of the old flat rounded picture.
"""

from pptx.enum.text import PP_ALIGN

from tools.pptx_core import (
    new_slide, add_meta, add_text, add_picture, add_rounded_rect, add_label_with_mark,
    FONT_REGULAR, FONT_MEDIUM, FONT_SEMIBOLD,
)


def build(prs, assets_dir):
    slide = new_slide(prs)
    add_meta(slide, "Bluewater Purifiers · 2026", "06/09")

    a = f"{assets_dir}/slide_06"

    # --- Left card: "Spirit purifier" intro ---
    add_rounded_rect(slide, 48, 112, 586, 526, "#f4f4f5", 32)
    add_text(slide, 96, 160, 490, 24, [[("Spirit purifier", FONT_MEDIUM, 24, "#2563eb", 0.9167)]])
    add_text(
        slide, 96, 362, 490, 106,
        [[("Small footprint. Powerful purification.", FONT_SEMIBOLD, 48, "#18181b", 1.1)]],
    )
    add_text(
        slide, 96, 484, 490, 106,
        [[("For homes, offices cafés and restaurants.", FONT_MEDIUM, 48, "#71717a", 1.1)]],
    )

    # --- Mini stat card: 3 L/min ---
    add_rounded_rect(slide, 48, 670, 586, 126, "#f4f4f5", 32)
    add_text(slide, 72, 712, 200, 41, [[("3 L/min", FONT_SEMIBOLD, 36, "#27272a", 1.15)]])
    add_text(
        slide, 329, 712, 281, 42,
        [[("Pure water flow rate. Spirit produces up to 3,800 L/day. No tank needed.", FONT_REGULAR, 16, "#27272a", 1.34)]],
        align=PP_ALIGN.RIGHT,
    )

    # --- Water recovery card: 70% ---
    add_rounded_rect(slide, 48, 828, 586, 204, "#f4f4f5", 32)
    add_picture(slide, f"{a}/fiber_thumb.png", 124, 828, 510, 204)
    add_text(slide, 72, 860, 337, 22, [[("Up to", FONT_REGULAR, 16, "#27272a", 1.375)]])
    add_text(slide, 72, 894, 337, 62, [[("70%", FONT_SEMIBOLD, 56, "#18181b", 1.107)]])
    add_text(
        slide, 72, 968, 271, 42,
        [
            [("Water recovery. Conventional ", FONT_REGULAR, 16, "#27272a", 1.34)],
            [("RO reaches ~25%.", FONT_REGULAR, 16, "#27272a", 1.34)],
        ],
    )

    # --- Center: product photo (background card + photo) ---
    add_rounded_rect(slide, 663, 112, 595.92, 920, "#f4f4f5", 32)
    add_picture(slide, f"{a}/spirit_photo.png", 663, 191, 596, 841)

    # --- Right: Certification card ---
    add_rounded_rect(slide, 1290, 112, 582, 206, "#f4f4f5", 32)
    add_picture(slide, f"{a}/nsf_badge.png", 1636, 155, 236, 120)
    add_text(slide, 1314, 144, 333, 22, [[("Certification", FONT_REGULAR, 16, "#27272a", 1.375)]])
    add_text(slide, 1314, 178, 333, 62, [[("NSF 372", FONT_SEMIBOLD, 56, "#18181b", 1.107)]])
    add_text(
        slide, 1314, 252, 354, 42,
        [[("NSF/ANSI/CAN 372. Product conforms with lead content requirements for “lead free” plumbing.", FONT_REGULAR, 16, "#27272a", 1.34)]],
    )

    # --- Right: SuperiorOsmosis card ---
    add_rounded_rect(slide, 1290, 350, 582, 682, "#f4f4f5", 32)
    add_picture(slide, f"{a}/osmosis_badge.png", 1290, 476, 582, 556)
    add_label_with_mark(
        slide, 1338, 398, "SuperiorOsmosis", FONT_SEMIBOLD, 36, "#18181b",
        "™", FONT_REGULAR, 23.22, w=486, h=40, line_spacing=1.1,
    )
    add_text(
        slide, 1338, 454, 486, 128,
        [[("Goes beyond ordinary filtration, using a self-cleaning membrane that removes up to 99.7% of contaminants for water that's pure and safe.", FONT_MEDIUM, 24, "#52525b", 1.34)]],
    )

    return slide
