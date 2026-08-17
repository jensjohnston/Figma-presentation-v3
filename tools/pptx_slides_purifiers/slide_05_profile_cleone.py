"""Slide 05 -- Cleone profile. Figma node 61219:18007.

Pixel-for-pixel the same content and geometry as the main deck's
slide_17_cleone.py (verified: every card position, image crop and copy
string matches, down to the exact clipped pixel dimensions of the "Reverse
osmosis" dome graphic and the WQA badge), so that slide's already-cropped
assets are reused directly.
"""

from pptx.enum.text import PP_ALIGN

from tools.pptx_core import (
    new_slide, add_meta, add_text, add_picture, add_rounded_rect,
    FONT_REGULAR, FONT_MEDIUM, FONT_SEMIBOLD,
)


def _mini_stat(slide, x, y, heading, heading_w, body, body_x, body_w):
    add_rounded_rect(slide, x, y, 586, 126, "#f4f4f5", 32)
    add_text(
        slide, x + 24, y + 42.5, heading_w, 41,
        [[(heading, FONT_SEMIBOLD, 36, "#27272a", 1.15)]],
    )
    add_text(
        slide, body_x, y + 44, body_w, 38,
        [[(body, FONT_REGULAR, 14, "#27272a", 1.34)]],
        align=PP_ALIGN.RIGHT,
    )


def build(prs, assets_dir):
    slide = new_slide(prs)
    add_meta(slide, "Bluewater Purifiers · 2026", "05/09")

    a = f"{assets_dir}/slide_05"

    # --- Left column: intro card ---
    add_rounded_rect(slide, 48, 112, 586.286, 444, "#f4f4f5", 32)
    add_text(
        slide, 96, 160, 490, 22,
        [[("Cleone purifier", FONT_MEDIUM, 24, "#2563eb", 1.0)]],
    )
    add_text(
        slide, 96, 280, 490, 106,
        [[("The smart choice for smaller-scale needs.", FONT_SEMIBOLD, 48, "#18181b", 1.1)]],
    )
    add_text(
        slide, 96, 402, 490, 106,
        [[("An affordable and reliable water purifier.", FONT_MEDIUM, 48, "#71717a", 1.1)]],
    )

    # --- Left column: mini stat cards ---
    _mini_stat(slide, 48, 588, "2.75 L/min", 178, "Flow from tank. Cleone produces up to 610 L/day.", 418, 192)
    _mini_stat(slide, 48, 747, "Tank-based", 346, "12 L tank included. The flow rate depends on tank level.", 434, 176)
    _mini_stat(slide, 48, 906, "Home & office", 239, "Cleone is perfect for homes and small offices.", 441, 169)

    # --- Center column: product photo card ---
    add_rounded_rect(slide, 663, 112, 595.92, 920, "#f4f4f5", 32)
    add_picture(slide, f"{a}/cleone_side.png", 663, 212, 596, 820)

    # --- Right column: Reverse osmosis card ---
    add_rounded_rect(slide, 1290, 112, 582, 444, "#f4f4f5", 32)
    add_picture(slide, f"{a}/osmosis_mesh.png", 1290, 334, 582, 222, radius_px=32)
    add_text(
        slide, 1338, 160, 486, 46,
        [[("Reverse osmosis", FONT_SEMIBOLD, 36, "#18181b", 1.1)]],
    )
    add_text(
        slide, 1338, 216, 486, 70,
        [[("Removes up to 99% of dissolved contaminants including heavy metals, bacteria, viruses, and more.", FONT_MEDIUM, 24, "#52525b", 1.34)]],
    )

    # --- Right column: Water recovery stat card ---
    add_rounded_rect(slide, 1290, 588, 582, 206, "#f4f4f5", 32)
    add_picture(slide, f"{a}/osmosis_zoom_thumb.png", 1366, 588, 506, 206, radius_px=32)
    add_text(
        slide, 1314, 620, 300, 22,
        [[("Up to", FONT_REGULAR, 16, "#27272a", 1.34)]],
    )
    add_text(
        slide, 1314, 654, 300, 62,
        [[("50%", FONT_SEMIBOLD, 56, "#18181b", 1.107)]],
    )
    add_text(
        slide, 1314, 728, 267, 42,
        [
            {"runs": [("Water recovery. Conventional", FONT_REGULAR, 16, "#27272a", 1.34)]},
            {"runs": [("RO reaches ~25%.", FONT_REGULAR, 16, "#27272a", 1.34)]},
        ],
    )

    # --- Right column: Certification card ---
    add_rounded_rect(slide, 1290, 826, 582, 206, "#f4f4f5", 32)
    add_picture(slide, f"{a}/wqa_badge.png", 1731, 874, 110, 110)
    add_text(
        slide, 1314, 858, 300, 22,
        [[("Certification", FONT_REGULAR, 16, "#27272a", 1.34)]],
    )
    add_text(
        slide, 1314, 892, 300, 62,
        [[("WQA", FONT_SEMIBOLD, 56, "#18181b", 1.107)]],
    )
    add_text(
        slide, 1314, 966, 281, 42,
        [[("NSF/ANSI/CAN 372: Drinking Water System Components - Lead Content.", FONT_REGULAR, 16, "#27272a", 1.34)]],
    )

    return slide
