"""Slide 08 -- Purifier Lineup. Figma node 61219:17997.

Background composition (Cleone + 2x Spirit close-ups) is pixel-identical to
slide_01_cover.py -- same node names, same object-fit percentages, same box
positions -- so its already-generated background PNGs are reused directly.
Only the headline text differs (a question, not the deck title) and there's
a standard meta-left/meta-right pair here (unlike the cover, which only had
a lone top-right line).

Pictures are drawn before text for the same reason as slide_01_cover.py:
their blank headspace was flattened to an opaque matte by get_screenshot,
so drawing them last would bury the headline/meta.
"""

from pptx.enum.text import PP_ALIGN

from tools.pptx_core import new_slide, add_meta, add_text, add_picture, add_rect, FONT_SEMIBOLD


def build(prs, assets_dir):
    slide = new_slide(prs)

    a = f"{assets_dir}/slide_08"

    add_rect(slide, 0, 0, 1920, 1080, "#f5f5f5")

    add_picture(slide, f"{a}/cleone_bg.png", 0, 0, 1561, 1080)
    add_picture(slide, f"{a}/spirit_bg_1.png", 782, 115, 1138, 965)
    add_picture(slide, f"{a}/spirit_bg_2.png", 0, 145, 1920, 935)

    add_meta(slide, "Bluewater Purifiers · 2026", "08/09")

    add_text(
        slide, 48, 115, 1805, 90,
        [[("Which Bluewater purifier is right for you?", FONT_SEMIBOLD, 80, "#27272a", 1.1)]],
        align=PP_ALIGN.CENTER,
    )

    return slide
