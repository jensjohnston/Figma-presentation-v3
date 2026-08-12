"""Slide 27 -- House Station. Figma node 63132:456."""

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from tools.pptx_core import new_slide, add_meta, add_text, add_picture, add_rounded_rect, add_label_with_mark, FONT_REGULAR, FONT_MEDIUM, FONT_SEMIBOLD


def _mini_stat_card(slide, x, y, heading, caption, caption_w):
    add_rounded_rect(slide, x, y, 586, 126, "#f4f4f5", 32)
    add_text(
        slide, x + 24, y, 300, 126,
        [[(heading, FONT_SEMIBOLD, 36, "#27272a", 1.15)]],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, x + 586 - 24 - caption_w, y, caption_w, 126,
        [[(caption, FONT_REGULAR, 14, "#27272a", 1.34)]],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )


def build(prs, assets_dir):
    slide = new_slide(prs)
    add_meta(slide, "Bluewater Stations", "27/33")

    # Column 1 top -- title card
    add_rounded_rect(slide, 48, 112, 586, 444, "#f4f4f5", 32)
    add_label_with_mark(
        slide, 96, 160, "House Station 1", FONT_MEDIUM, 24, "#2563eb",
        "™", FONT_REGULAR, 12, w=400, h=30, line_spacing=1.1,
    )
    add_text(
        slide, 96, 281, 490, 230,
        [
            {"runs": [("Pure water. Every tap. Every appliance.", FONT_SEMIBOLD, 48, "#18181b", 1.1)], "space_after_px": 16},
            {"runs": [("Custom-made whole-house solution.", FONT_MEDIUM, 48, "#71717a", 1.1)]},
        ],
    )

    # Column 1 bottom -- 3 mini stat cards
    _mini_stat_card(slide, 48, 588, "7,600 L/day", "Powered by Pro purifier. Up to 7,600 L/day of pure drinking water.", 219)
    _mini_stat_card(slide, 48, 747, "Water sources", "Purifies municipal, well, lake, and brackish water.", 178)
    _mini_stat_card(slide, 48, 906, "Tailor-made", "Custom-configured to match your water source, household size, and daily needs.", 262)

    # Column 2 -- product photo (flattened, already rounded + clipped by Figma)
    add_picture(slide, f"{assets_dir}/slide_27/product.png", 666, 112, 588, 920)

    # Column 3 -- "99.7% pure water" card
    add_rounded_rect(slide, 1286, 112, 586, 920, "#f4f4f5", 32)
    add_text(
        slide, 1334, 160, 490, 50,
        [[("99.7% pure water", FONT_SEMIBOLD, 36, "#27272a", 1.34)]],
    )
    add_text(
        slide, 1334, 224, 490, 135,
        [
            [
                ("Powered by Pro purifier and Bluewater SuperiorOsmosis", FONT_MEDIUM, 24, "#52525b", 1.34),
                ("™", FONT_MEDIUM, 24, "#52525b", 1.34),
                (", House Station 1 removes up to 99.7% of contaminants for water that's pure, safe, and trusted.", FONT_MEDIUM, 24, "#52525b", 1.34),
            ]
        ],
    )
    # Decorative glow sphere + "SuperiorOsmosis" gradient wordmark -- flattened as one
    # image (gradient-filled text on a multi-layer blurred glow is impractical to
    # reproduce as editable pptx text/shapes).
    add_picture(slide, f"{assets_dir}/slide_27/glow.png", 1286, 372, 586, 660)

    return slide
