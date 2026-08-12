"""Slide 06 -- St Andrews: The Home of Golf. Figma node 63222:131.

Full-bleed photo (with its baked-in dark overlay) + white centered wordmark logo +
large white headline. Background and logo are screenshotted directly since they carry
no separately-reconstructed text (the "ST ANDREWS" wordmark inside the logo asset is
part of the club's crest/logo lockup, a decorative brand mark).
"""

from pptx.enum.text import PP_ALIGN

from tools.pptx_core import new_slide, add_meta, add_text, add_picture, FONT_SEMIBOLD


def build(prs, assets_dir):
    slide = new_slide(prs)

    add_picture(slide, f"{assets_dir}/slide_06/bg_photo.png", 0, 0, 1920, 1080)
    add_picture(slide, f"{assets_dir}/slide_06/logo.png", 862, 316, 197, 75)

    add_text(
        slide, 48, 415, 1824, 240,
        [
            [("The home of golf", FONT_SEMIBOLD, 120, "#ffffff", 1.0)],
            [("runs on Bluewater.", FONT_SEMIBOLD, 120, "#ffffff", 1.0)],
        ],
        align=PP_ALIGN.CENTER,
    )

    add_meta(slide, "Client Story", "06/33")
    return slide
