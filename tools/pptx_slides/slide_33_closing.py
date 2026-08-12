"""Slide 33 -- Closing. Figma node 63150:136.

No meta-left on this slide (the Bluewater logo takes that spot instead), so this
does not use add_meta() -- meta-right is added directly to match its exact style.
"""

from pptx.enum.text import PP_ALIGN

from tools.pptx_core import new_slide, add_text, add_picture, FONT_REGULAR, FONT_SEMIBOLD


def build(prs, assets_dir):
    slide = new_slide(prs)

    add_text(slide, 1600, 48, 272, 19, [[("33/33", FONT_REGULAR, 14, "#71717a", 1.0)]], align=PP_ALIGN.RIGHT)

    # Logo (icon + wordmark) -- flattened, screenshot of the whole logo instance.
    add_picture(slide, f"{assets_dir}/slide_33/logo.png", 48, 46, 334, 64)

    # Headline. Figma uses a blue linear-gradient text fill (#3b82f6 -> #1d4ed8); pptx text
    # runs only support solid color, so this is approximated with brand blue-600 (#2563eb),
    # the midpoint of that gradient (see core/brand-tokens.json).
    add_text(
        slide, 48, 403, 953, 274,
        [
            [("Water like", FONT_SEMIBOLD, 152, "#2563eb", 0.9)],
            [("never before.", FONT_SEMIBOLD, 152, "#2563eb", 0.9)],
        ],
    )

    # Hero photo (hand holding bottle) -- flattened, screenshot of the exact photo node.
    add_picture(slide, f"{assets_dir}/slide_33/photo.png", 938, 88, 983, 992)

    return slide
