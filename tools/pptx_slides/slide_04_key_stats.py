"""Slide 04 -- Key Stats. Figma node 63145:142.

A dense 3-column bento grid. Judgment calls:
- Two text elements use a Figma linear-gradient blue fill ("Water like never before."
  and "99.7% / pure water."); approximated with solid #2c68e7 (average of the gradient
  stops), same choice as slide 01.
- "For the planet." heading uses a custom script font ("Bluescript:Bold") for the word
  "planet." only, not available in this deck's font set; falls back to FONT_SEMIBOLD
  (plain, no script styling) for the whole line.
- The blue "SuperiorOsmosis" circle graphic, the ingredients/bottles background photos,
  and the heart/drop/recycle icons are all purely decorative -> flattened images,
  screenshotted directly from Figma (already correctly clipped to their visible
  on-canvas bounds by get_screenshot's default ancestor-clip behavior).
- "plastic bottles" is genuine strikethrough content -> real text with a manual
  OOXML `strike` attribute (python-pptx has no high-level strikethrough API).
"""

from pptx.oxml.ns import qn
from pptx.enum.text import MSO_ANCHOR

from tools.pptx_core import new_slide, add_meta, add_text, add_picture, add_rounded_rect, FONT_MEDIUM, FONT_SEMIBOLD

BLUE = "#2c68e7"


def _strike(text_box):
    for para in text_box.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set("strike", "sngStrike")


def build(prs, assets_dir):
    slide = new_slide(prs)
    add_meta(slide, "Bluewater", "04/33")
    a = assets_dir

    # ---- Column 1 (x=48, w=587) ----
    add_rounded_rect(slide, 48, 115, 587, 284, "#f4f4f5", 32)
    add_text(slide, 96, 163, 490, 25, [[("Bluewater", FONT_MEDIUM, 16, "#3f3f46", 1.34)]])
    add_text(
        slide, 96, 200, 490, 151,
        [
            [("Water like ", FONT_SEMIBOLD, 80, BLUE, 1.0)],
            [("never before.", FONT_SEMIBOLD, 80, BLUE, 1.0)],
        ],
        anchor=MSO_ANCHOR.BOTTOM,
    )

    add_rounded_rect(slide, 48, 431, 587, 284, "#f4f4f5", 32)
    add_picture(slide, f"{a}/slide_04/purification_graphic.png", 220, 431, 415, 284)
    add_text(slide, 79, 516, 460, 30, [[("Purification", FONT_MEDIUM, 16, "#3f3f46", 1.34)]])
    add_text(
        slide, 79, 544, 460, 100,
        [
            {"runs": [("99.7% ", FONT_SEMIBOLD, 48, BLUE, 1.1)]},
            {"runs": [("pure water.", FONT_MEDIUM, 28, BLUE, 1.15)]},
        ],
    )

    add_picture(slide, f"{a}/slide_04/ingredients_bg.png", 48, 747, 277, 284, radius_px=32)
    add_text(slide, 72, 771, 229, 25, [[("Ingredients", FONT_MEDIUM, 16, "#3f3f46", 1.1)]])
    add_text(slide, 72, 798, 229, 60, [[("Healthy ingredients.", FONT_SEMIBOLD, 24, "#18181b", 1.34)]])

    add_rounded_rect(slide, 357, 747, 277, 284, "#f4f4f5", 32)
    add_picture(slide, f"{a}/slide_04/bottles_bg.png", 357, 830, 277, 201)
    add_text(slide, 381, 771, 229, 25, [[("Bottles", FONT_MEDIUM, 16, "#3f3f46", 1.1)]])
    add_text(slide, 381, 798, 229, 60, [[("Reusable bottles.", FONT_SEMIBOLD, 24, "#18181b", 1.34)]])

    # ---- Column 2 (x=666, w=587) -- hero photo card, no text ----
    add_picture(slide, f"{a}/slide_04/card_hero.png", 666, 115, 587, 915)

    # ---- Column 3 (x=1286, w=587) ----
    add_rounded_rect(slide, 1286, 115, 587, 221, "#f4f4f5", 32)
    add_text(
        slide, 1318, 147, 523, 157,
        [
            [("For you.", FONT_SEMIBOLD, 48, "#18181b", 1.1)],
            [("For people.", FONT_SEMIBOLD, 48, "#18181b", 1.1)],
            [("For the planet.", FONT_SEMIBOLD, 48, "#18181b", 1.1)],
        ],
    )

    add_rounded_rect(slide, 1286, 368, 277, 339, "#f4f4f5", 32)
    add_text(slide, 1310, 392, 229, 25, [[("For you", FONT_MEDIUM, 16, "#3f3f46", 1.1)]])
    add_text(slide, 1310, 429, 229, 90, [[("You're choosing health over harmful ingredients.", FONT_SEMIBOLD, 24, "#18181b", 1.34)]])
    add_picture(slide, f"{a}/slide_04/heart_icon.png", 1453, 601, 86, 85)

    add_rounded_rect(slide, 1595, 368, 277, 339, "#f4f4f5", 32)
    add_text(slide, 1619, 392, 229, 25, [[("For people", FONT_MEDIUM, 16, "#3f3f46", 1.1)]])
    add_text(slide, 1619, 426, 229, 90, [[("You're supporting projects delivering clean, safe water.", FONT_SEMIBOLD, 24, "#18181b", 1.34)]])
    add_picture(slide, f"{a}/slide_04/drop_icon.png", 1772, 600, 86, 86)

    add_rounded_rect(slide, 1286, 739, 587, 291, "#f4f4f5", 32)
    add_text(slide, 1318, 771, 400, 25, [[("For the planet", FONT_MEDIUM, 16, "#3f3f46", 1.34)]])
    strike_box = add_text(slide, 1318, 814, 389, 65, [[("plastic bottles", FONT_SEMIBOLD, 56, "#17a34a", 1.0)]])
    _strike(strike_box)
    add_text(slide, 1318, 900, 304, 70, [[("You are saying no to single-use plastics and throwaway choices.", FONT_MEDIUM, 24, "#71717a", 1.34)]])
    add_picture(slide, f"{a}/slide_04/recycle_icon.png", 1733, 857, 166, 188)

    return slide
