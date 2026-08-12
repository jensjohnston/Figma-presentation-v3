"""Slide 18 -- Spirit purifier. Figma node 65080:4042."""

from tools.pptx_core import (
    new_slide, add_meta, add_text, add_picture, add_rounded_rect,
    FONT_REGULAR, FONT_MEDIUM, FONT_SEMIBOLD,
)
from pptx.enum.text import PP_ALIGN


def build(prs, assets_dir):
    slide = new_slide(prs)
    add_meta(slide, "Bluewater Purifiers", "18/33")

    a = f"{assets_dir}/slide_18"

    # --- Left card: "Spirit purifier" intro ---
    add_rounded_rect(slide, 48, 112, 586, 526, "#f4f4f5", 32)
    add_text(slide, 96, 160, 490, 24, [[("Spirit purifier", FONT_MEDIUM, 24, "#2563eb", 0.9167)]])
    add_text(
        slide, 96, 362, 490, 106,
        [[("Small footprint. Powerful purification.", FONT_SEMIBOLD, 48, "#18181b", 1.1)]],
    )
    add_text(
        slide, 96, 484, 490, 106,
        [[("For homes, offices cafés and restaurants.", FONT_MEDIUM, 48, "#71717a", 1.1)]],
    )

    # --- Mini stat card: 3 L/min ---
    add_rounded_rect(slide, 48, 670, 586, 126, "#f4f4f5", 32)
    add_text(slide, 72, 712, 200, 41, [[("3 L/min", FONT_SEMIBOLD, 36, "#27272a", 1.15)]])
    add_text(
        slide, 329, 712, 281, 42,
        [[("Pure water flow rate. Spirit produces up to 3,800 L/day. No tank needed.", FONT_REGULAR, 16, "#27272a", 1.34)]],
        align=PP_ALIGN.RIGHT,
    )

    # --- Water recovery card: 70% ---
    add_rounded_rect(slide, 48, 828, 586, 204, "#f4f4f5", 32)
    add_picture(slide, f"{a}/fiber_thumb.png", 124, 828, 510, 204)
    add_text(slide, 72, 860, 337, 22, [[("Up to", FONT_REGULAR, 16, "#27272a", 1.375)]])
    add_text(slide, 72, 894, 337, 62, [[("70%", FONT_SEMIBOLD, 56, "#18181b", 1.107)]])
    add_text(
        slide, 72, 968, 271, 42,
        [
            [("Water recovery. Conventional ", FONT_REGULAR, 16, "#27272a", 1.34)],
            [("RO reaches ~25%.", FONT_REGULAR, 16, "#27272a", 1.34)],
        ],
    )

    # --- Center: product photo ---
    add_picture(slide, f"{a}/photo.png", 663, 191, 596, 841, radius_px=32)

    # --- Right: Certification card ---
    add_rounded_rect(slide, 1290, 112, 582, 206, "#f4f4f5", 32)
    add_picture(slide, f"{a}/nsf_badge.png", 1659, 155, 236, 120)
    add_text(slide, 1314, 144, 333, 22, [[("Certification", FONT_REGULAR, 16, "#27272a", 1.375)]])
    add_text(slide, 1314, 178, 333, 62, [[("NSF 372", FONT_SEMIBOLD, 56, "#18181b", 1.107)]])
    add_text(
        slide, 1314, 252, 354, 42,
        [[("NSF/ANSI/CAN 372. Product conforms with lead content requirements for “lead free” plumbing.", FONT_REGULAR, 16, "#27272a", 1.34)]],
    )

    # --- Right: SuperiorOsmosis card ---
    add_rounded_rect(slide, 1290, 350, 582, 682, "#f4f4f5", 32)
    add_picture(slide, f"{a}/osmosis_badge.png", 1290, 476, 582, 556)
    add_text(
        slide, 1338, 398, 486, 40,
        [
            [
                ("SuperiorOsmosis", FONT_SEMIBOLD, 36, "#18181b", 1.1),
                ("™", FONT_REGULAR, 23.22, "#18181b", 1.1),
            ]
        ],
    )
    add_text(
        slide, 1338, 454, 486, 128,
        [[("Goes beyond ordinary filtration, using a self-cleaning membrane that removes up to 99.7% of contaminants for water that's pure and safe.", FONT_MEDIUM, 24, "#52525b", 1.34)]],
    )

    return slide
