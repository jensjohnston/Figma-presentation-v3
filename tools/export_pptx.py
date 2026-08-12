"""
Reconstructs Figma slides as an editable PowerPoint deck.

Text layers become real, editable PowerPoint text boxes (font/size/color/position
matched to the Figma design). Everything else (photos, icons, decorative shapes,
rounded cards) is rendered as flattened images exported directly from Figma via
get_screenshot on the specific non-text node -- this avoids reverse-engineering
Figma's CSS transforms (object-cover crops, mirrors, rotations) and guarantees
pixel-perfect backgrounds.

Slide content is hand-authored per slide (positions/text pulled from Figma's
get_design_context output) into the SLIDES list below. This is a pilot covering
two representative slides; the same build_slide() machinery is meant to scale to
the full deck.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# Figma canvas is 1920x1080px, treated as 144 "dpi" so 1920px -> 13.333in (16:9 widescreen).
PX = Emu(6350)  # EMU per Figma px
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080

FONT_REGULAR = "Suisse Int'l"
FONT_MEDIUM = "Suisse Int'l Medium"
FONT_SEMIBOLD = "Suisse Int'l Semi Bold"


def px(v):
    return Emu(round(v * 6350))


def pt(figma_px):
    # 144 "dpi" canvas -> 2 Figma px per point
    return Pt(figma_px / 2)


def rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_corner_radius(shape, radius_px, box_w_px, box_h_px):
    """Apply a rounded-rectangle clip to a picture or autoshape via its adj value."""
    adj = min(0.5, radius_px / min(box_w_px, box_h_px))
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    geom = spPr.find(qn("a:prstGeom"))
    if geom is None:
        # picture placeholders use a:prstGeom too once we set it
        for tag in ("a:custGeom", "a:prstGeom"):
            existing = spPr.find(qn(tag))
            if existing is not None:
                spPr.remove(existing)
        geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
        spPr.append(geom)
    else:
        geom.set("prst", "roundRect")
        for child in list(geom):
            geom.remove(child)
    avLst = geom.makeelement(qn("a:avLst"), {})
    gd = avLst.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {int(adj * 100000)}"})
    avLst.append(gd)
    geom.append(avLst)


def add_picture(slide, path, x_px, y_px, w_px, h_px, radius_px=0):
    pic = slide.shapes.add_picture(path, px(x_px), px(y_px), px(w_px), px(h_px))
    if radius_px:
        set_corner_radius(pic, radius_px, w_px, h_px)
    return pic


def add_rounded_rect(slide, x_px, y_px, w_px, h_px, fill_hex, radius_px):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x_px), px(y_px), px(w_px), px(h_px))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.fill.background()
    shape.shadow.inherit = False
    set_corner_radius(shape, radius_px, w_px, h_px)
    return shape


def add_text(slide, x_px, y_px, w_px, h_px, paragraphs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of paragraph specs, each a list of run tuples
    (text, font_name, size_px, color_hex, line_spacing) or a dict with 'runs' + 'space_after_px'."""
    box = slide.shapes.add_textbox(px(x_px), px(y_px), px(w_px), px(h_px))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, para_spec in enumerate(paragraphs):
        if isinstance(para_spec, dict):
            runs = para_spec["runs"]
            space_after = para_spec.get("space_after_px", 0)
        else:
            runs = para_spec
            space_after = 0

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        line_spacing = runs[0][4] if runs else None
        size_px = runs[0][2] if runs else None
        if line_spacing and size_px:
            # Figma's line-height is an exact multiple of font size (e.g. 1.34 * 32px).
            # A bare float here sets PowerPoint's *percentage* spacing mode, which is a
            # multiple of its own "single line spacing" (font-metrics based, taller than
            # 1.0x the font size) -- not the font size itself, so it renders too tall.
            # Pt() switches to exact-spacing mode, matching Figma's definition precisely.
            p.line_spacing = Pt((size_px * line_spacing) / 2)
        if space_after:
            p.space_after = pt(space_after)
        else:
            p.space_after = Pt(0)
        p.space_before = Pt(0)

        for text, font_name, size_px, color_hex, _ls in runs:
            r = p.add_run()
            r.text = text
            r.font.name = font_name
            r.font.size = pt(size_px)
            r.font.color.rgb = rgb(color_hex)
    return box


def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    return slide


def add_meta(slide, left_text, right_text):
    add_text(slide, 48, 48, 300, 19, [[(left_text, FONT_REGULAR, 14, "#71717a", 1.0)]])
    add_text(slide, 1600, 48, 272, 19, [[(right_text, FONT_REGULAR, 14, "#71717a", 1.0)]], align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 02 -- Who We Are
# ---------------------------------------------------------------------------

def build_slide_02(prs, assets_dir):
    slide = new_slide(prs)
    add_meta(slide, "Bluewater", "02/33")

    add_text(
        slide, 48, 115, 896, 140,
        [
            [("Meet Bengt Rittri. ", FONT_SEMIBOLD, 64, "#18181b", 1.1)],
            [("Founder of Bluewater.", FONT_SEMIBOLD, 64, "#18181b", 1.1)],
        ],
    )

    body_runs_1 = [
        ("Growing up on the beaches of southern Sweden, Bengt Rittri spent his childhood searching for amber. But over the years, ", FONT_MEDIUM, 32, "#71717a", 1.34),
        ("the amber gave way to plastic bottles", FONT_MEDIUM, 32, "#3f3f46", 1.34),
        (" — and ", FONT_MEDIUM, 32, "#71717a", 1.34),
        ("he knew something had to change.", FONT_MEDIUM, 32, "#3f3f46", 1.34),
    ]
    body_runs_2 = [
        ("In 2013, Rittri founded Bluewater.", FONT_MEDIUM, 32, "#3f3f46", 1.34),
        (" At Bluewater, we set out to reimagine what drinking water could be: using ", FONT_MEDIUM, 32, "#71717a", 1.34),
        ("advanced purification and natural ingredients", FONT_MEDIUM, 32, "#3f3f46", 1.34),
        (" to create the healthiest, most sustainable choice ", FONT_MEDIUM, 32, "#71717a", 1.34),
        ("for you and the world around you.", FONT_MEDIUM, 32, "#3f3f46", 1.34),
    ]
    add_text(
        slide, 48, 602, 850, 430,
        [
            {"runs": body_runs_1, "space_after_px": 32},
            {"runs": body_runs_2},
        ],
    )

    add_picture(slide, f"{assets_dir}/slide02_photo_card.png", 976, 115, 896, 917)
    return slide


# ---------------------------------------------------------------------------
# Slide 32 -- Customise Your Bottle
# ---------------------------------------------------------------------------

def bento_cell_text(slide, x, y, label, heading, body):
    add_text(
        slide, x, y, 792, 190,
        [
            {"runs": [(label, FONT_MEDIUM, 24, "#52525b", 1.34)], "space_after_px": 16},
            {"runs": [(heading, FONT_SEMIBOLD, 36, "#27272a", 1.15)], "space_after_px": 16},
            {"runs": [(body, FONT_MEDIUM, 24, "#52525b", 1.34)]},
        ],
    )


def build_slide_32(prs, assets_dir):
    slide = new_slide(prs)
    add_meta(slide, "Bluewater Bottles", "32/33")

    add_text(
        slide, 48, 115, 1400, 90,
        [[("Make it yours. Customize your bottle.", FONT_SEMIBOLD, 64, "#212126", 1.0)]],
    )

    # Left card: Silicone loop
    add_rounded_rect(slide, 48, 287, 888, 745, "#f4f4f5", 32)
    add_picture(slide, f"{assets_dir}/slide32_left_photo.png", 48, 550, 888, 482, radius_px=32)
    add_picture(slide, f"{assets_dir}/swatches_left.png", 96, 488, 330, 24)
    bento_cell_text(slide, 96, 335, "Silicone loop", "Carry your bottle with ease", "Durable silicone loop. Available in 10 colors.")

    # Right card: Sports cap
    add_rounded_rect(slide, 984, 287, 888, 745, "#f4f4f5", 32)
    add_picture(slide, f"{assets_dir}/slide32_right_photo.png", 984, 550, 888, 482, radius_px=32)
    add_picture(slide, f"{assets_dir}/swatches_right.png", 1032, 488, 126, 24)
    bento_cell_text(slide, 1032, 335, "Sports cap", "Lid designed for active people", "Smart sports-cap. Available in 4 colors.")

    return slide


def main():
    import sys
    assets_dir = sys.argv[1] if len(sys.argv) > 1 else "/tmp/pptx_assets"
    out_path = sys.argv[2] if len(sys.argv) > 2 else "/tmp/bluewater_pilot.pptx"

    prs = Presentation()
    prs.slide_width = px(SLIDE_W_PX)
    prs.slide_height = px(SLIDE_H_PX)

    build_slide_02(prs, assets_dir)
    build_slide_32(prs, assets_dir)

    prs.save(out_path)
    print(f"Saved {out_path}")


if __name__ == "__main__":
    main()
